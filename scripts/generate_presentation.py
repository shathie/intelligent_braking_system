"""
Generate a PPTX dissertation presentation for:
Intelligent Multi-Modal Braking System for Friction-Aware Adaptive Braking Control
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x47, 0xA1)   # title backgrounds
MID_BLUE    = RGBColor(0x15, 0x65, 0xC0)   # accents / section headers
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)   # slide body background
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT      = RGBColor(0xFF, 0x6F, 0x00)   # orange for highlights
GREEN       = RGBColor(0x1B, 0x5E, 0x20)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)


def _rgb(r, g, b):
    return RGBColor(r, g, b)


# ── Helper: fill shape solid colour ─────────────────────────────────────────
def solid_fill(shape, rgb):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


# ── Helper: add a text frame with a single paragraph ─────────────────────────
def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


# ── Helper: banner across the top of a content slide ───────────────────────
def add_banner(slide, title_text, slide_width=Inches(13.33),
               banner_height=Inches(1.15)):
    banner = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, slide_width, banner_height)
    solid_fill(banner, DARK_BLUE)
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


# ── Helper: slide background fill ───────────────────────────────────────────
def set_bg(slide, rgb=LIGHT_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


# ── Helper: bullet list ──────────────────────────────────────────────────────
def add_bullets(slide, left, top, width, height, bullets,
                font_size=16, color=DARK_TEXT, indent=False):
    """bullets: list of str, or list of (str, level) tuples"""
    from pptx.util import Pt, Inches
    from pptx.oxml.ns import qn
    import copy, lxml.etree as etree

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = level
        run = p.add_run()
        bullet_char = "•  " if level == 0 else "    ◦  "
        run.text = bullet_char + text
        run.font.size = Pt(font_size - level * 1.5)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(4)

    return txBox


# ── Helper: two-column layout  ───────────────────────────────────────────────
def two_col(slide, left_bullets, right_bullets,
            left_top=1.3, right_top=1.3,
            col_w=5.9, gap=0.3, font_size=15):
    add_bullets(slide, 0.3, left_top, col_w, 5.5, left_bullets,
                font_size=font_size)
    add_bullets(slide, 0.3 + col_w + gap, right_top, col_w, 5.5,
                right_bullets, font_size=font_size)


# ── Helper: coloured info box ────────────────────────────────────────────────
def info_box(slide, left, top, width, height, label, value,
             bg=MID_BLUE, label_color=WHITE, val_color=WHITE,
             label_size=13, val_size=22):
    box = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    solid_fill(box, bg)
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = label
    r0.font.size = Pt(label_size)
    r0.font.color.rgb = label_color
    r0.font.name = "Calibri"

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(val_size)
    r1.font.bold = True
    r1.font.color.rgb = val_color
    r1.font.name = "Calibri"


# ════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank

# ── SLIDE 1 · Title ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BLUE)

# Gradient-like accent bar at top
bar = s.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(0.18))
solid_fill(bar, ACCENT)
bar.line.fill.background()

# Institution label
add_textbox(s, 0.3, 0.25, 12.7, 0.5,
            "Birla Institute of Technology & Science, Pilani — WILP Division",
            font_size=14, color=_rgb(0xBB, 0xDE, 0xFF), align=PP_ALIGN.CENTER)

# Main title
title_box = s.shapes.add_textbox(
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Intelligent Multi-Modal Braking System\nfor Friction-Aware Adaptive Braking Control"
r.font.size = Pt(38)
r.font.bold = True
r.font.color.rgb = WHITE
r.font.name = "Calibri"

# Subtitle
add_textbox(s, 0.5, 3.85, 12.3, 0.6,
            "WILP Dissertation Presentation  |  August 2026",
            font_size=18, color=_rgb(0xBB, 0xDE, 0xFF), align=PP_ALIGN.CENTER)

# Divider
div = s.shapes.add_shape(1, Inches(3.5), Inches(4.55), Inches(6.3), Inches(0.04))
solid_fill(div, ACCENT)
div.line.fill.background()

# Keywords row
add_textbox(s, 0.5, 4.7, 12.3, 0.55,
            "Computer Vision  ·  Physics-Informed ML  ·  Reinforcement Learning  ·  Adaptive Control  ·  Intelligent Transportation",
            font_size=14, color=_rgb(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

add_textbox(s, 0.5, 6.7, 12.3, 0.5,
            "BITS Pilani  |  Intelligent Braking System  |  Dissertation 2026",
            font_size=11, color=_rgb(0x78, 0xAB, 0xD4), align=PP_ALIGN.CENTER)


# ── SLIDE 2 · Agenda ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Agenda")

agenda_items = [
    "1.  Introduction & Motivation",
    "2.  Literature Review",
    "3.  Research Methodology & System Architecture",
    "4.  Data Collection & Preprocessing",
    "5.  Experimental Results & Analysis",
    "6.  Conclusion & Future Work",
    "7.  Demo",
]

for i, item in enumerate(agenda_items):
    top = 1.35 + i * 0.78
    box = s.shapes.add_shape(
        1, Inches(1.5), Inches(top), Inches(10.3), Inches(0.62))
    bg_col = MID_BLUE if i % 2 == 0 else _rgb(0x1E, 0x88, 0xE5)
    solid_fill(box, bg_col)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "   " + item
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"


# ── SLIDE 3 · Introduction — Problem Statement ───────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Introduction  —  The Problem")

add_textbox(s, 0.3, 1.25, 12.7, 0.45,
            "Why does braking performance degrade under varying road conditions?",
            font_size=17, bold=True, color=MID_BLUE)

left_b = [
    "Tire-road friction coefficient (μ) varies widely:",
    ("Dry asphalt  →  μ ≈ 0.8", 1),
    ("Wet road      →  μ ≈ 0.5", 1),
    ("Icy surface   →  μ ≈ 0.1", 1),
    "Conventional braking systems assume fixed or simplified friction",
    "Fixed μ assumptions cause:",
    ("Extended stopping distances", 1),
    ("Wheel lock on slippery surfaces", 1),
    ("Increased accident risk in adverse weather", 1),
]

right_b = [
    "Real-world challenge:",
    ("Visual cues alone are ambiguous (wet vs. shadowed asphalt)", 1),
    ("Vehicle signals alone miss environmental context", 1),
    ("Physics must be respected for safe estimation", 1),
    "",
    "This project addresses it by building an intelligent,",
    "multi-modal, friction-aware braking system end-to-end.",
]

two_col(s, left_b, right_b, font_size=14)


# ── SLIDE 4 · Introduction — Objectives ─────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Introduction  —  Objectives")

objectives = [
    ("Build a reproducible end-to-end pipeline: data → train → evaluate → simulate → report", 0),
    ("Learn visual road-surface representations using a Vision Transformer (ViT)", 0),
    ("Model temporal vehicle dynamics from CAN-like signal sequences", 0),
    ("Fuse visual and temporal modalities to reduce prediction ambiguity", 0),
    ("Constrain friction estimation with physics via a Physics-Informed Neural Network (PINN)", 0),
    ("Learn an adaptive braking policy using Soft Actor-Critic (SAC) reinforcement learning", 0),
    ("Produce measurable, academically defensible baseline and improved metrics", 0),
]

add_bullets(s, 0.5, 1.3, 12.3, 5.8, objectives, font_size=16)


# ── SLIDE 5 · Literature Review ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Literature Review")

lit = [
    "Vision Transformers (ViT)  [Dosovitskiy et al., ICLR 2021]",
    ("Self-attention over image patches captures richer scene semantics than CNNs for road classification", 1),
    "Physics-Informed Neural Networks (PINN)  [Raissi et al., J. Comp. Physics 2019]",
    ("Composite loss = data loss + physics residual enforces vehicle-dynamics consistency", 1),
    ("Prevents numerically plausible but physically impossible friction outputs", 1),
    "Soft Actor-Critic (SAC)  [Haarnoja et al., ICML 2018]",
    ("Maximum-entropy RL objective → sample-efficient, stable stochastic braking policies", 1),
    "Multimodal Fusion in Autonomous Driving",
    ("Cross-modal attention merges camera and signal streams, outperforming single-modality baselines", 1),
    "Gap addressed by this work:",
    ("No prior system integrates ViT + temporal LSTM/GRU + cross-modal fusion + PINN + SAC"
     " in one reproducible braking pipeline", 1),
]

add_bullets(s, 0.4, 1.3, 12.5, 5.8, lit, font_size=14)


# ── SLIDE 6 · Research Methodology — Architecture ────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Research Methodology  —  System Architecture")

add_textbox(s, 0.3, 1.2, 12.7, 0.4,
            "Five-stage perception-estimation-control pipeline",
            font_size=15, bold=True, color=MID_BLUE)

# Pipeline boxes
stages = [
    ("1\nViT\nClassifier", _rgb(0x01, 0x57, 0x9B)),
    ("2\nTemporal\nNetwork", _rgb(0x00, 0x69, 0x9A)),
    ("3\nFusion\nNetwork", _rgb(0x00, 0x7B, 0x83)),
    ("4\nPINN\nFriction", _rgb(0x00, 0x89, 0x53)),
    ("5\nSAC\nControl", _rgb(0x33, 0x69, 0x1E)),
]

box_w, box_h = 2.1, 1.7
gap = 0.15
start_x = 0.35
top_y = 1.75

for i, (label, col) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    b = s.shapes.add_shape(1, Inches(x), Inches(top_y), Inches(box_w), Inches(box_h))
    solid_fill(b, col)
    b.line.fill.background()
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    # Arrow (except after last)
    if i < len(stages) - 1:
        ax = x + box_w + 0.01
        arr = s.shapes.add_shape(1, Inches(ax), Inches(top_y + 0.7),
                                 Inches(gap + 0.0), Inches(0.3))
        solid_fill(arr, _rgb(0xFF, 0x6F, 0x00))
        arr.line.fill.background()

# Input labels
inputs = ["Road image\n(camera)", "CAN signals\n(speed, accel…)", "ViT + Temporal\nfeatures", "Fused\nlatent vec", "μ estimate\n+ state"]
for i, inp in enumerate(inputs):
    x = start_x + i * (box_w + gap)
    add_textbox(s, x, top_y + box_h + 0.05, box_w, 0.6, inp,
                font_size=10, color=_rgb(0x0D, 0x47, 0xA1), align=PP_ALIGN.CENTER)

# Architecture notes
notes = [
    "Training order is intentional: upstream encoders are frozen before downstream blocks are trained",
    "Cross-modal attention in the Fusion Network resolves modality ambiguity",
    "PINN composite loss:  L_total = α·L_data + β·L_physics   (warm-up scheduling applied)",
    "SAC learns stochastic braking policy from friction-aware state using maximum-entropy objective",
    "Reproducible workflow orchestrated through main.py with CLI stage control",
]
add_bullets(s, 0.3, 3.75, 12.7, 3.4, notes, font_size=13)


# ── SLIDE 7 · Research Methodology — Design Decisions ───────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Research Methodology  —  Key Design Decisions")

left_b = [
    "Why ViT instead of CNN?",
    ("Self-attention over 16×16 patches captures global context", 1),
    ("Pre-trained vit_base_patch16_224 → 768-dim embeddings", 1),
    "",
    "Why multimodal fusion?",
    ("Dark asphalt = dry shadow OR wet surface → ambiguous visually", 1),
    ("Temporal deceleration/wheel-speed trends disambiguate", 1),
    ("Cross-modal attention fuses both into shared 512-dim space", 1),
    "",
    "Why PINN?",
    ("Data-only regression can violate braking physics", 1),
    ("Physics residual: F_braking = μ·m·g (wheel-force equation)", 1),
]

right_b = [
    "Why SAC for control?",
    ("Off-policy, sample-efficient, handles continuous action spaces", 1),
    ("Maximum-entropy term prevents policy collapse", 1),
    ("Actor LR: 3e-4 | Critic LR: 1e-3 | γ: 0.99", 1),
    "",
    "Why CPU-only execution?",
    ("Quadro T1200 WDDM (4 GB) OOM at batch_size=32 with ViT", 1),
    ("CPU gave equivalent throughput (~9-15 s/batch) without OOM", 1),
    "",
    "Reproducibility controls:",
    ("YAML config files for all hyperparameters", 1),
    ("Environment-variable overrides for runtime caps", 1),
    ("Timestamped backups before each pipeline run", 1),
]

two_col(s, left_b, right_b, font_size=13)


# ── SLIDE 8 · Data Collection ────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Data Collection  —  Datasets")

datasets = [
    ("THU Road Surface", "958,941", "Road surface texture classification (27 classes)"),
    ("Mendeley Vehicle",  "370,151", "Multi-modal: images + CAN-like signals"),
    ("DAWN Weather",       "1,027",  "Adverse-weather road scenes (fog, snow, rain, sand)"),
    ("BDD100K (subset)", " 26,000", "Diverse driving imagery across conditions"),
    ("KITTI Raw (subset)","    910", "Driving sequences with stereo cameras"),
]

header_cols = ["Dataset", "Samples", "Role"]
col_ws = [3.5, 1.6, 7.3]
col_xs = [0.3, 3.85, 5.5]
header_y = 1.25

# header row
for ci, (hdr, cw, cx) in enumerate(zip(header_cols, col_ws, col_xs)):
    hb = s.shapes.add_shape(1, Inches(cx), Inches(header_y),
                             Inches(cw), Inches(0.42))
    solid_fill(hb, DARK_BLUE)
    hb.line.fill.background()
    tf = hb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = hdr
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

row_colors = [_rgb(0xE3, 0xF2, 0xFD), WHITE]
for ri, (name, count, role) in enumerate(datasets):
    ry = header_y + 0.42 + ri * 0.55
    row_bg = row_colors[ri % 2]
    for ci, (val, cw, cx) in enumerate(
            zip([name, count, role], col_ws, col_xs)):
        rb = s.shapes.add_shape(1, Inches(cx), Inches(ry),
                                 Inches(cw), Inches(0.52))
        solid_fill(rb, row_bg)
        rb.line.color.rgb = _rgb(0xBB, 0xDE, 0xFF)
        rb.line.width = Pt(0.5)
        tf = rb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci < 2 else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ("  " if ci == 2 else "") + val
        r.font.size = Pt(13)
        r.font.color.rgb = DARK_TEXT
        r.font.name = "Calibri"
        if ci == 0:
            r.font.bold = True

add_textbox(s, 0.3, 4.4, 12.7, 0.4,
            "Total: ~1.36 M samples across 5 sources  |  Dataset analysis automated via scripts/analyze_datasets.py",
            font_size=13, bold=False, color=MID_BLUE)

notes2 = [
    "Class imbalance identified: THU and BDD100K dominate; DAWN and KITTI underrepresented",
    "Runtime sample caps (env vars BDD_MAX_SAMPLES, KITTI_MAX_SAMPLES) used for practical experimentation",
]
add_bullets(s, 0.3, 4.95, 12.7, 1.5, notes2, font_size=13)


# ── SLIDE 9 · Data Collection — Preprocessing ────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Data Collection  —  Preprocessing Pipeline")

steps = [
    ("Raw Input", "Heterogeneous: JPEG, PNG, raw KITTI bins, CSV", _rgb(0x45, 0x6A, 0xAE)),
    ("Image Preprocessing", "Resize → 224×224, normalize (ImageNet mean/std), tensor", _rgb(0x00, 0x7B, 0x9E)),
    ("CAN Synthesis", "Synthetic speed, accel, jerk, wheel-speed, brake signals (17-dim)", _rgb(0x00, 0x89, 0x60)),
    ("Target Generation", "μ labels mapped from surface class via friction lookup table", _rgb(0x33, 0x69, 0x1E)),
    ("Output Artifacts", "image tensors (.pt), CAN tensors (.pt), targets.pt, metadata.csv", _rgb(0x6A, 0x1B, 0x9A)),
]

for i, (title, desc, col) in enumerate(steps):
    y = 1.3 + i * 1.05
    lb = s.shapes.add_shape(1, Inches(0.3), Inches(y), Inches(2.5), Inches(0.85))
    solid_fill(lb, col)
    lb.line.fill.background()
    tf = lb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    desc_b = s.shapes.add_shape(1, Inches(2.9), Inches(y), Inches(10.1), Inches(0.85))
    solid_fill(desc_b, _rgb(0xF5, 0xF9, 0xFF))
    desc_b.line.color.rgb = col
    desc_b.line.width = Pt(1)
    tf2 = desc_b.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "  " + desc
    r2.font.size = Pt(14)
    r2.font.color.rgb = DARK_TEXT
    r2.font.name = "Calibri"

    # connector dot
    dot = s.shapes.add_shape(9,  # oval
                              Inches(2.7), Inches(y + 0.3),
                              Inches(0.15), Inches(0.25))
    solid_fill(dot, col)
    dot.line.fill.background()


# ── SLIDE 10 · Results — Classification (ViT) ────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Results  —  Road Surface Classification  (ViT)")

# Metric boxes
metrics = [
    ("Accuracy", "69.27%", MID_BLUE),
    ("Macro F1", "0.4233", _rgb(0x1B, 0x5E, 0x20)),
    ("Weighted F1", "0.7019", _rgb(0x4A, 0x14, 0x8C)),
    ("Macro Precision", "0.4393", _rgb(0x9E, 0x46, 0x00)),
]

for i, (lbl, val, col) in enumerate(metrics):
    x = 0.3 + i * 3.2
    info_box(s, x, 1.25, 2.9, 1.5, lbl, val, bg=col)

# Comparison vs baseline
add_textbox(s, 0.3, 2.95, 12.7, 0.35,
            "Progressive improvement over project lifecycle:",
            font_size=14, bold=True, color=DARK_BLUE)

compare = [
    ("Stage", "Accuracy", "Macro F1", "Weighted F1"),
    ("Early Baseline", "16.41%", "0.0616", "0.1448"),
    ("Mid-Semester Run", "37.50%", "0.1574", "0.3197"),
    ("Latest (Final) Run", "69.27%", "0.4233", "0.7019"),
]

header_bg = [DARK_BLUE, _rgb(0xFF, 0xC1, 0x07), _rgb(0xFF, 0xC1, 0x07), _rgb(0xFF, 0xC1, 0x07)]
row_bgs   = [_rgb(0xE3, 0xF2, 0xFD), WHITE, _rgb(0xE8, 0xF5, 0xE9)]
cws = [3.8, 2.5, 2.5, 2.8]
cxs = [0.3, 4.15, 6.7, 9.25]

for ri, row in enumerate(compare):
    ry = 3.45 + ri * 0.62
    for ci, (val, cw, cx) in enumerate(zip(row, cws, cxs)):
        rb = s.shapes.add_shape(1, Inches(cx), Inches(ry), Inches(cw), Inches(0.55))
        if ri == 0:
            solid_fill(rb, DARK_BLUE)
            txt_col = WHITE
        elif ri == 3:
            solid_fill(rb, _rgb(0xC8, 0xE6, 0xC9))
            txt_col = _rgb(0x1B, 0x5E, 0x20)
        else:
            solid_fill(rb, row_bgs[ri-1])
            txt_col = DARK_TEXT
        rb.line.color.rgb = _rgb(0xBB, 0xDE, 0xFF)
        rb.line.width = Pt(0.5)
        tf = rb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        r.font.size = Pt(13)
        r.font.bold = (ri == 0 or ri == 3)
        r.font.color.rgb = txt_col if ri > 0 else WHITE
        r.font.name = "Calibri"

add_textbox(s, 0.3, 6.78, 12.7, 0.35,
            "Macro vs. Weighted F1 gap indicates persisting class imbalance — addressed in recommendations",
            font_size=12, color=_rgb(0x6A, 0x1B, 0x9A))


# ── SLIDE 11 · Results — Regression (Temporal / Fusion / PINN) ───────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Results  —  Friction Estimation  (Temporal · Fusion · PINN)")

add_textbox(s, 0.3, 1.2, 12.7, 0.35,
            "All three regression stages output μ (friction coefficient) on the same evaluation split",
            font_size=14, bold=False, color=MID_BLUE)

reg_data = [
    ("Model", "MSE", "RMSE", "MAE", "R²", "Max Error"),
    ("Temporal Network", "≈ 0.0000", "≈ 0.0000", "3.5e-41", "1.0000", "1.4e-40"),
    ("Fusion Network",   "0.0191",   "0.1382",   "0.1000",  "0.5863", "0.6129"),
    ("PINN (Final)",     "0.0187",   "0.1368",   "0.0962",  "0.5931", "0.6449"),
]

cws2 = [3.0, 1.8, 1.8, 1.8, 1.8, 2.1]
cxs2 = [0.3, 3.35, 5.2, 7.05, 8.9, 10.75]

for ri, row in enumerate(reg_data):
    ry = 1.7 + ri * 0.7
    for ci, (val, cw, cx) in enumerate(zip(row, cws2, cxs2)):
        rb = s.shapes.add_shape(1, Inches(cx), Inches(ry), Inches(cw), Inches(0.62))
        if ri == 0:
            solid_fill(rb, DARK_BLUE); txt_col = WHITE
        elif ri == 3:
            solid_fill(rb, _rgb(0xC8, 0xE6, 0xC9)); txt_col = _rgb(0x1B, 0x5E, 0x20)
        else:
            solid_fill(rb, WHITE if ri % 2 == 0 else _rgb(0xE3, 0xF2, 0xFD))
            txt_col = DARK_TEXT
        rb.line.color.rgb = _rgb(0xBB, 0xDE, 0xFF)
        rb.line.width = Pt(0.5)
        tf = rb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        r.font.size = Pt(12)
        r.font.bold = (ri == 0 or ri == 3)
        r.font.color.rgb = txt_col if ri > 0 else WHITE
        r.font.name = "Calibri"

interpret = [
    "Temporal Network R²=1.0: perfect reconstruction loss — stage completion confirmed (not direct braking evidence)",
    "PINN (RMSE 0.1368, R²=0.5931) slightly outperforms Fusion — physics regularisation is beneficial",
    "PINN physics loss:  L = α·MSE(μ_pred, μ_true) + β·||F_brake − μ·m·g||  constrains output to vehicle dynamics",
    "Positive R² on PINN confirms practically useful friction predictions under current data distribution",
]
add_bullets(s, 0.3, 4.6, 12.7, 2.7, interpret, font_size=13, color=DARK_TEXT)


# ── SLIDE 12 · Results — Control (SAC) ──────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Results  —  Adaptive Braking Control  (SAC)")

ctrl_data = [
    ("Surface", "Stop Distance (m)", "Stop Time (s)", "Max Jerk (m/s³)", "Stop Success", "Timeouts"),
    ("Dry",     "25.1",              "2.5",            "0.96",             "100%",         "0/10"),
    ("Wet",     "50.2",              "5.0",            "0.16",             "100%",         "0/10"),
    ("Icy",     "194.1",             "19.0",           "0.01",             "100%",         "0/10"),
    ("Rough",   "33.5",              "3.3",            "0.46",             "100%",         "0/10"),
]

cws3 = [2.0, 2.5, 2.1, 2.1, 2.0, 1.8]
cxs3 = [0.2, 2.25, 4.8, 6.95, 9.1, 11.15]

for ri, row in enumerate(ctrl_data):
    ry = 1.3 + ri * 0.68
    for ci, (val, cw, cx) in enumerate(zip(row, cws3, cxs3)):
        rb = s.shapes.add_shape(1, Inches(cx), Inches(ry), Inches(cw), Inches(0.6))
        if ri == 0:
            solid_fill(rb, DARK_BLUE); txt_col = WHITE
        else:
            solid_fill(rb, WHITE if ri % 2 == 0 else _rgb(0xE3, 0xF2, 0xFD))
            txt_col = DARK_TEXT
            if ci == 4:  # success column — highlight 100%
                solid_fill(rb, _rgb(0xC8, 0xE6, 0xC9)); txt_col = _rgb(0x1B, 0x5E, 0x20)
        rb.line.color.rgb = _rgb(0xBB, 0xDE, 0xFF)
        rb.line.width = Pt(0.5)
        tf = rb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        r.font.size = Pt(12)
        r.font.bold = (ri == 0)
        r.font.color.rgb = txt_col if ri > 0 else WHITE
        r.font.name = "Calibri"

# Highlight box
hb = s.shapes.add_shape(1, Inches(0.2), Inches(4.45), Inches(12.9), Inches(0.62))
solid_fill(hb, _rgb(0xFF, 0xF9, 0xC4))
hb.line.color.rgb = ACCENT
hb.line.width = Pt(1.5)
tf = hb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "  ✔ 100% stop success on all surfaces  |  Icy 194 m (μ=0.1) vs Dry 25 m (μ=0.8) — 7.7× ratio matches friction physics"
r.font.size = Pt(13)
r.font.bold = True
r.font.color.rgb = _rgb(0x6D, 0x40, 0x00)
r.font.name = "Calibri"

ctrl_notes = [
    "Simulator bug fixed: brake-torque sign error caused braking to accelerate the vehicle (T_net = +T_brake → -T_brake)",
    "SAC reward redesigned: 7× stronger velocity urgency penalty + time-scaled stop bonus (50–100 range)",
    "Retrained 50 epochs on corrected simulator; best epoch reward −38.3 ≈ theoretical optimum −33",
    "Low jerk across all surfaces (≤ 0.96 m/s³) — smooth deceleration consistent with passenger comfort",
]
add_bullets(s, 0.3, 5.2, 12.7, 2.1, ctrl_notes, font_size=13, color=DARK_TEXT)


# ── SLIDE 13 · Conclusion ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Conclusion")

left_b = [
    "What was delivered:",
    ("End-to-end intelligent braking pipeline — fully reproducible", 1),
    ("ViT classifier: 69.27% accuracy (↑ from 16.4% baseline)", 1),
    ("PINN friction estimator: RMSE 0.1368, R² 0.5931", 1),
    ("SAC control: 100% stop success on all 4 surfaces", 1),
    ("Stopping distances: Dry 25 m · Wet 50 m · Rough 34 m · Icy 194 m", 1),
    "",
    "Key engineering findings:",
    ("Identified & fixed brake-torque sign bug in vehicle simulator", 1),
    ("Redesigned SAC reward with urgency signal + time-scaled stop bonus", 1),
    ("Stabilised PINN with warm-up + physics residual normalisation", 1),
    ("Windows-safe multiprocessing for all DataLoaders", 1),
]

right_b = [
    "Limitations identified:",
    ("Class imbalance: macro F1 < weighted F1 across all stages", 1),
    ("CPU-only execution limits scalability", 1),
    ("SAC policy trained on simplified simulator — real-sensor noise not modelled", 1),
    "",
    "Recommendations:",
    ("Class-balanced sampling + minority-class augmentation", 1),
    ("Stage-specific epoch schedules (not global override)", 1),
    ("Ablation study: ViT vs. Temporal vs. Fusion vs. PINN", 1),
    ("Extend evaluation to intermediate μ values and sensor noise injection", 1),
    ("GPU-optimised DataLoader with preprocessed .pt files", 1),
]

two_col(s, left_b, right_b, font_size=13)


# ── SLIDE 14 · Demo Overview ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, LIGHT_BLUE)
add_banner(s, "Demo  —  System Walkthrough")

demo_steps = [
    ("1. Pipeline Entry Point", "python main.py --all",
     "Orchestrates all stages: analyze → preprocess → train → evaluate → simulate → report",
     _rgb(0x01, 0x57, 0x9B)),
    ("2. Training", "python scripts/train.py",
     "Trains ViT → Temporal → Fusion → PINN → SAC sequentially with stage isolation",
     _rgb(0x00, 0x69, 0x9A)),
    ("3. Evaluation", "python scripts/evaluate.py",
     "Generates per-model JSON metrics: vit_metrics.json, pinn_metrics.json, sac_metrics.json …",
     _rgb(0x00, 0x7B, 0x83)),
    ("4. Simulation", "python scripts/simulate.py",
     "Runs SAC braking policy on 4 surface types; produces stopping_distance.png, trajectories",
     _rgb(0x00, 0x89, 0x53)),
    ("5. Report Generation", "python scripts/report.py",
     "Aggregates metrics + plots into evaluation_report.html + dissertation draft",
     _rgb(0x33, 0x69, 0x1E)),
]

for i, (title, cmd, desc, col) in enumerate(demo_steps):
    y = 1.3 + i * 1.12
    # Title label
    tb = s.shapes.add_shape(1, Inches(0.2), Inches(y), Inches(2.7), Inches(0.95))
    solid_fill(tb, col)
    tb.line.fill.background()
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    # Command box
    cb = s.shapes.add_shape(1, Inches(2.98), Inches(y), Inches(2.8), Inches(0.95))
    solid_fill(cb, _rgb(0x1A, 0x1A, 0x2E))
    cb.line.fill.background()
    tf2 = cb.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = cmd
    r2.font.size = Pt(11)
    r2.font.bold = True
    r2.font.color.rgb = _rgb(0x80, 0xFF, 0xBD)
    r2.font.name = "Courier New"

    # Description
    db = s.shapes.add_shape(1, Inches(5.85), Inches(y), Inches(7.2), Inches(0.95))
    solid_fill(db, _rgb(0xF5, 0xF9, 0xFF))
    db.line.color.rgb = col
    db.line.width = Pt(1)
    tf3 = db.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run()
    r3.text = "  " + desc
    r3.font.size = Pt(12)
    r3.font.color.rgb = DARK_TEXT
    r3.font.name = "Calibri"


# ── SLIDE 15 · Thank You ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BLUE)

bar2 = s.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(0.18))
solid_fill(bar2, ACCENT)
bar2.line.fill.background()

add_textbox(s, 0.5, 1.8, 12.3, 1.4,
            "Thank You",
            font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

div2 = s.shapes.add_shape(1, Inches(3.5), Inches(3.4), Inches(6.3), Inches(0.05))
solid_fill(div2, ACCENT)
div2.line.fill.background()

add_textbox(s, 0.5, 3.6, 12.3, 0.6,
            "Intelligent Multi-Modal Braking System for Friction-Aware Adaptive Braking Control",
            font_size=17, color=_rgb(0xBB, 0xDE, 0xFF), align=PP_ALIGN.CENTER)

add_textbox(s, 0.5, 4.35, 12.3, 0.5,
            "BITS Pilani  |  WILP Division  |  Dissertation 2026",
            font_size=14, color=_rgb(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

add_textbox(s, 0.5, 5.2, 12.3, 1.5,
            "Questions Welcome",
            font_size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ── SAVE ─────────────────────────────────────────────────────────────────────
out_path = os.path.join(
    os.path.dirname(__file__), "..", "output", "reports",
    "dissertation_presentation.pptx")
out_path = os.path.normpath(out_path)
prs.save(out_path)
print(f"Saved: {out_path}")
